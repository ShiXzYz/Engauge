import os
from pdfminer.high_level import extract_text as extract_pdf_text
from pptx import Presentation


def extract_text_from_file(path: str) -> str:
    """Extract text from PDF or PPTX file path. Returns concatenated text."""
    if not os.path.exists(path):
        return ""
    lower = path.lower()
    if lower.endswith('.pdf'):
        try:
            return extract_pdf_text(path)
        except Exception:
            return ''
    if lower.endswith('.pptx') or lower.endswith('.ppt'):
        try:
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        texts.append(shape.text)
            return '\n'.join(texts)
        except Exception:
            return ''
    # fallback: try to read as text
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception:
        return ''
